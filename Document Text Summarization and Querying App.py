import streamlit as st
from docx import Document
import pptx
import os
from io import BytesIO
import PyPDF2
import nltk
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from transformers import pipeline, AutoModelForQuestionAnswering
import sys

# Function to read txt file
def read_txt(file):
    text = file.getvalue().decode("utf-8")
    return text

# Function to read docx file
def read_docx(file):
    doc = Document(BytesIO(file.read()))
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return '\n'.join(text)

# Function to read pptx file
def read_pptx(file):
    prs = pptx.Presentation(BytesIO(file.read()))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

# Function to read PDF file
def read_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text += page.extract_text()
    # Split the text into lines based on newlines
    text_lines = text.split("\n")
    # Join the lines back together with spaces
    text = " ".join(text_lines)
    return text

def summarize_text(input_text):
    # Initialize BART tokenizer and model
    checkpoint = "sshleifer/distilbart-cnn-12-6"
    tokenizer = AutoTokenizer.from_pretrained(checkpoint)
    model = AutoModelForSeq2SeqLM.from_pretrained(checkpoint)
    nltk.download('punkt')

    input_text = input_text.strip()
    input_text = input_text.strip("\r")
    
    # Convert file content to sentences
    sentences = nltk.tokenize.sent_tokenize(input_text)

    # Initialize chunks
    sum_of_max_tokens = sum(sorted([len(tokenizer.tokenize(sentence)) for sentence in sentences], reverse=True)[:3])

    # Create the chunks
    length = 0
    chunk = ""
    chunks = []
    for sentence in sentences:
        combined_length = len(tokenizer.tokenize(sentence)) + length
        if combined_length <= sum_of_max_tokens:
            chunk += sentence + " "
            length = combined_length
            if sentence == sentences[-1]:
                chunks.append(chunk.strip())
        else:
            chunks.append(chunk.strip())
            length = 0
            chunk = sentence + " "
            length = len(tokenizer.tokenize(sentence))

    # Generate summary
    summarized_text = ""
    for chunk in chunks:
        input_data = tokenizer(chunk, return_tensors="pt", max_length=int(0.40 * len(input_text)), truncation=True)
        output = model.generate(**input_data)
        summarized_text += tokenizer.decode(output[0], skip_special_tokens=True) + " "
    return summarized_text

# Add a button to download the summary
def download_summary_button(summarized_text):
    # Split the summary into lines
    summary_lines = summarized_text.split('. ')
    # Encode each line separately and join them with newline characters
    encoded_lines = [line.encode('utf-8') for line in summary_lines]
    data = b'. '.join(encoded_lines)
    # Add a button to download the summary
    st.download_button(
        label="Download Summary",
        data=data,
        file_name="summary.txt",
        mime="text/plain"
    )

# Function to query text using Hugging Face's question-answering pipeline
def query_text(input_text, question):
    # Load the model and tokenizer
    model_name ="deepset/roberta-base-squad2"

    model = AutoModelForQuestionAnswering.from_pretrained(model_name)
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    
    # Initialize the question-answering pipeline
    qa_pipeline = pipeline("question-answering", model=model, tokenizer=tokenizer)
    
    # Split the input text into chunks
    max_chunk_length = 600  # Maximum length for each chunk
    chunks = [input_text[i:i+max_chunk_length] for i in range(0, len(input_text), max_chunk_length)]
    
    # Fine-tune the model on each chunk and generate the answer
    final_answer = ""
    for chunk in chunks:
        # Generate the answer to the question for the current chunk
        result = qa_pipeline(question=question, context=chunk)
        answer = result["answer"]
        
        # Add the answer to the final answer string
        final_answer += answer + " "
    
    # Capitalize the first letter and add a period at the end
    final_answer = final_answer.strip().capitalize() + "."
    
    return final_answer

def main():
    st.set_page_config(layout="wide")  # Set page layout to wide
    st.title("Document Text Summarization & Querying App ✧⁠◝⁠(⁠⁰⁠▿⁠⁰⁠)⁠◜⁠✧")
    file = st.file_uploader("Upload File:", type=['txt', 'docx', 'pptx', 'pdf'])
    user_input = st.text_area("Enter Text:")
    # Partition interface into two columns
    col1, col2 = st.columns([1,1])
    with col1:
        with st.expander("Click to display text."):
            input_text = ""
            if file is not None:
                # Display uploaded file
                file_type = os.path.splitext(file.name)[1]
                if file_type == '.txt':
                    input_text = read_txt(file)
                elif file_type == '.docx':
                    input_text = read_docx(file)
                elif file_type == '.pptx':
                    input_text = read_pptx(file)
                elif file_type == '.pdf':
                    input_text = read_pdf(file)
                else:
                    st.write("Unsupported file format!")
                st.text_area("Make any necessary changes:", input_text, height=600)
                
            elif user_input != "":
                # If direct text input is provided
                input_text = user_input
                st.text_area("Make any necessary changes:", input_text, height=600)

    with col2:
        if input_text and isinstance(input_text, str):
            with st.expander("Click to select an action."):
                # Select box for further actions
                action = st.selectbox("✧⁠◝⁠(⁠⁰⁠▿⁠⁰⁠)⁠◜⁠✧",["Select an option", "Summarize Text", "Query Text", "Exit App"])

                if action == "Summarize Text":
                    summarized_text = summarize_text(input_text)
                    st.text_area("Summarized Text", summarized_text, height=500)
                    download_summary_button(summarized_text)
                elif action == "Query Text":
                    question = st.text_input("Enter your question:")
                    if st.button("Query") :
                        if question.strip():  # Check if the question is not empty
                            answer = query_text(input_text, question)
                            st.write("Answer:", answer)
                        else:
                            st.write("Please enter a question.")
                elif action == "Exit App":
                    st.warning("Are you sure you want to exit?")
                    if st.button("Yes"):
                        os._exit(0)
                elif action =="Select an option":
                    pass

if __name__ == "__main__":
    main()
